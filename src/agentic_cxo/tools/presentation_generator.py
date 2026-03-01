"""
Presentation Generator Tool — create modern PPT slides per brand guidelines.

Generates PowerPoint presentations with content from research or user input,
applying company brand colors, fonts, and visual style from BrandStore.
Available to any CXO for marketing decks, investor pitches, internal briefings.
"""

from __future__ import annotations

import re
import uuid
import logging
from pathlib import Path
from typing import Any

from agentic_cxo.tools.brand_intelligence import BrandStore, BrandProfile
from agentic_cxo.tools.framework import BaseTool, ToolParam, ToolResult

logger = logging.getLogger(__name__)

DATA_DIR = Path(".cxo_data") / "presentations"
DEFAULT_COLORS = ["#6366f1", "#8b5cf6", "#4f46e5", "#7c3aed", "#312e81"]


def _hex_to_rgb(hex_val: str) -> tuple[int, int, int]:
    """Convert #RRGGBB to (r, g, b) 0-255."""
    hex_val = hex_val.strip().lstrip("#")
    if len(hex_val) == 3:
        hex_val = "".join(c * 2 for c in hex_val)
    if len(hex_val) != 6:
        return (99, 102, 241)  # Default indigo
    return (
        int(hex_val[0:2], 16),
        int(hex_val[2:4], 16),
        int(hex_val[4:6], 16),
    )


def _parse_outline(text: str) -> list[dict[str, Any]]:
    """Parse user outline into slide structure."""
    slides: list[dict[str, Any]] = []
    lines = text.strip().split("\n")
    current_slide: dict[str, Any] | None = None

    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Slide title: ## or ### or "Slide N:" or leading number
        if line.startswith("##") or re.match(r"^#+\s+.+", line):
            if current_slide:
                slides.append(current_slide)
            title = re.sub(r"^#+\s*", "", line).strip()
            current_slide = {"title": title, "bullets": []}
        elif re.match(r"^(\d+)[.)]\s+", line):
            if current_slide is None:
                current_slide = {"title": "Content", "bullets": []}
            bullet = re.sub(r"^\d+[.)]\s*", "", line).strip()
            if bullet:
                current_slide["bullets"].append(bullet)
        elif line.startswith("-") or line.startswith("*"):
            if current_slide is None:
                current_slide = {"title": "Content", "bullets": []}
            bullet = line.lstrip("-* ").strip()
            if bullet:
                current_slide["bullets"].append(bullet)
        else:
            if current_slide is None:
                current_slide = {"title": line[:80], "bullets": []}
            else:
                current_slide["bullets"].append(line)

    if current_slide:
        slides.append(current_slide)

    if not slides:
        slides = [{"title": "Untitled", "bullets": []}]
    return slides


class PresentationGeneratorTool(BaseTool):
    """Create PowerPoint presentations with brand guideline compliance."""

    def __init__(self) -> None:
        self._store = BrandStore()

    @property
    def name(self) -> str:
        return "presentation_generator"

    @property
    def description(self) -> str:
        return (
            "Create a PowerPoint (.pptx) presentation with modern slides. "
            "Uses company brand guidelines (colors, fonts) from BrandStore "
            "when available. Provide a title and outline or bullet points. "
            "Use for: marketing decks, investor pitches, internal briefings, "
            "sales presentations, product overviews."
        )

    @property
    def parameters(self) -> list[ToolParam]:
        return [
            ToolParam(
                name="title",
                description="Presentation title",
            ),
            ToolParam(
                name="outline",
                description=(
                    "Slide outline: use ## for slide titles, - or 1. for bullets. "
                    "Example: ## Intro\\n- Point 1\\n- Point 2\\n## Next Slide\\n- Item"
                ),
            ),
            ToolParam(
                name="brand_domain",
                description="Optional: company domain for brand colors (e.g. acme.com)",
                required=False,
            ),
        ]

    @property
    def trigger_keywords(self) -> list[str]:
        return [
            "create presentation", "make ppt", "powerpoint", "slide deck",
            "create slides", "presentation deck", "pitch deck",
            "marketing deck", "investor presentation", "create a deck",
            "ppt slides", "presentation slides", "branded presentation",
        ]

    def execute(
        self,
        title: str = "",
        outline: str = "",
        brand_domain: str = "",
        **kwargs: Any,
    ) -> ToolResult:
        if not title and not outline:
            return ToolResult(
                tool_name=self.name,
                success=False,
                error="Provide at least a title or outline for the presentation.",
            )

        title = (title or "Presentation").strip()
        if not outline:
            outline = f"## {title}\n- Key points to cover\n\n## Next Steps\n- Action items"

        brand = self._get_brand(brand_domain)
        slides_data = _parse_outline(outline)

        try:
            pptx_path = self._generate_pptx(title, slides_data, brand)
        except ImportError as e:
            return ToolResult(
                tool_name=self.name,
                success=False,
                error=f"python-pptx not installed: {e}. Run: pip install python-pptx",
            )
        except Exception as e:
            logger.exception("PPT generation failed")
            return ToolResult(
                tool_name=self.name,
                success=False,
                error=str(e),
            )

        static_path = self._copy_to_static(pptx_path)
        report = (
            f"## Presentation Created: {title}\n\n"
            f"**{len(slides_data)} slides** generated with "
            f"{'brand guidelines' if brand else 'default styling'}.\n\n"
            f"📄 **[Download PowerPoint]({static_path})**"
        )

        return ToolResult(
            tool_name=self.name,
            success=True,
            data={
                "title": title,
                "slides_count": len(slides_data),
                "path": str(pptx_path),
                "url": static_path,
                "brand_used": brand.company_name or brand.domain if brand else None,
            },
            summary=report,
        )

    def _get_brand(self, domain: str | None) -> BrandProfile | None:
        if domain and domain.strip():
            return self._store.get(domain.strip().replace("www.", ""))
        return self._store.primary_brand

    def _generate_pptx(
        self,
        title: str,
        slides_data: list[dict[str, Any]],
        brand: BrandProfile | None,
    ) -> Path:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        colors = self._get_color_palette(brand)
        title_rgb = _hex_to_rgb(colors[0])
        accent_rgb = _hex_to_rgb(colors[1]) if len(colors) > 1 else title_rgb

        blank_layout = prs.slide_layouts[6]  # Blank

        for i, slide_info in enumerate(slides_data):
            slide_title = slide_info.get("title", "Slide")
            bullets = slide_info.get("bullets", [])

            slide = prs.slides.add_slide(blank_layout)

            # Background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

            # Title shape
            left, top, width, height = Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
            title_box = slide.shapes.add_textbox(left, top, width, height)
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*title_rgb)
            p.font.name = (brand.heading_font or "Calibri") if brand else "Calibri"

            # Bullets
            if bullets:
                bullet_left = Inches(0.5)
                bullet_top = Inches(1.4)
                bullet_width = Inches(12.333)
                bullet_height = Inches(5.5)
                bullet_box = slide.shapes.add_textbox(
                    bullet_left, bullet_top, bullet_width, bullet_height
                )
                bf = bullet_box.text_frame
                bf.word_wrap = True
                for j, bullet in enumerate(bullets[:8]):
                    if j == 0:
                        p = bf.paragraphs[0]
                    else:
                        p = bf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(60, 60, 60)
                    p.font.name = (brand.body_font or "Calibri") if brand else "Calibri"
                    p.space_after = Pt(8)

            # Accent bar at bottom
            line_shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), Inches(7.0), Inches(12.333), Inches(0.05),
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(*accent_rgb)
            line_shape.line.fill.background()

        DATA_DIR.mkdir(parents=True, exist_ok=True)
        filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
        path = DATA_DIR / filename
        prs.save(str(path))
        return path

    def _get_color_palette(self, brand: BrandProfile | None) -> list[str]:
        if brand and brand.all_colors:
            return brand.all_colors[:5]
        if brand and brand.primary_color:
            return [
                brand.primary_color,
                brand.secondary_color or brand.primary_color,
                brand.accent_color or brand.primary_color,
            ]
        return DEFAULT_COLORS

    def _copy_to_static(self, pptx_path: Path) -> str:
        import shutil

        static_dir = Path("src/agentic_cxo/api/static/presentations")
        static_dir.mkdir(parents=True, exist_ok=True)
        dest = static_dir / pptx_path.name
        shutil.copy2(str(pptx_path), str(dest))
        return f"/static/presentations/{pptx_path.name}"
