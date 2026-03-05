"""
Premium presentation generator — McKinsey/BCG-grade slide design.

Features:
- Multi-tone title typography (gold highlight on key terms)
- Section category labels on every slide
- Rich layout types: data_metrics, research_study, definition_boxes,
  two_column_info, warning_callout, benefits_risks, comparison_table,
  recommendations, bottom_line, sources, quote
- Bar charts via manual rectangles
- Consistent dark-navy + gold premium brand palette
"""

from __future__ import annotations

import re
import uuid
import logging
import datetime as dt
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

logger = logging.getLogger(__name__)

_PKG_ROOT = Path(__file__).resolve().parent.parent.parent.parent
DATA_DIR = (_PKG_ROOT / ".cxo_data" / "presentations").resolve()
if not _PKG_ROOT.exists() or not (_PKG_ROOT / "src").exists():
    DATA_DIR = Path(".cxo_data").resolve() / "presentations"


# ── Colour helpers ──────────────────────────────────────────────

def _hex(h: str) -> RGBColor:
    h = h.strip().lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        h = "0F172A"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _clean(t: str) -> str:
    t = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", t)
    t = re.sub(r"\*\*(.+?)\*\*", r"\1", t)
    t = re.sub(r"\*(.+?)\*", r"\1", t)
    t = re.sub(r"`(.+?)`", r"\1", t)
    return t.strip()


def _bullets(text: str) -> list[str]:
    out: list[str] = []
    parent_bullet: str = ""
    for line in text.split("\n"):
        stripped = line.strip()
        indent = len(line) - len(line.lstrip())
        m = re.match(r"^[-*•]\s+(.+)", stripped) or re.match(r"^\d+[.)]\s+(.+)", stripped)
        if m:
            content = _clean(m.group(1).strip())
            if indent >= 4 and parent_bullet:
                out.append(f"{parent_bullet}: {content}")
            else:
                parent_bullet = content
                out.append(content)
        elif stripped.startswith(">"):
            quote_text = _clean(stripped.lstrip("> ").strip())
            if quote_text:
                out.append(quote_text)
        elif re.match(r"^[A-Z].*:\s+.+", stripped):
            out.append(_clean(stripped))
    return out


def _parse_sections(content: str) -> list[dict[str, Any]]:
    sections: list[dict[str, Any]] = []
    title = ""
    body_lines: list[str] = []
    for line in content.split("\n"):
        m = re.match(r"^#{1,3}\s+(.+)$", line)
        if m:
            if title or body_lines:
                body = "\n".join(body_lines).strip()
                sections.append({"title": title, "body": body, "bullets": _bullets(body)})
            title = m.group(1).strip()
            body_lines = []
        else:
            body_lines.append(line)
    if title or body_lines:
        body = "\n".join(body_lines).strip()
        sections.append({"title": title, "body": body, "bullets": _bullets(body)})
    return sections


# ── Shape / text helpers ────────────────────────────────────────

def _rect(slide, left, top, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def _textbox(slide, left, top, w, h, text, size=14, color=None, bold=False,
             align=PP_ALIGN.LEFT, font="Calibri", line_spacing=None, space_after=0):
    if not text:
        return None
    color = color or RGBColor(0xFF, 0xFF, 0xFF)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _clean(str(text))
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    if space_after:
        p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return box


def _rich_title(slide, left, top, w, h, title, h_font, highlight_color,
                base_color=None, base_size=32, hl_size=52, after_size=26):
    """Multi-run textbox: intro text (base) + key noun phrase (gold/large) + tail (gray).

    Works for ANY topic by splitting at a natural break point — prepositions/conjunctions
    ('of', 'on', 'in', 'and', 'the') mark the boundary between intro words and the
    key concept phrase. Falls back to splitting after the first 2 words.
    """
    base_color = base_color or RGBColor(0xFF, 0xFF, 0xFF)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True

    words = title.split()

    def _find_split(words: list[str]) -> tuple[str, str, str]:
        """Return (before, highlight, after). Highlight = the dominant noun phrase."""
        # Strategy 1: split at a preposition/article that follows 1-3 intro words
        split_words = {"of", "on", "in", "at", "for", "and", "the", "to", "by", "with", "about"}
        for i in range(1, min(4, len(words))):
            if words[i].lower() in split_words:
                # before: words 0..i-1, highlight: words i+1..end (skip the preposition)
                before  = " ".join(words[:i])
                hl_part = " ".join(words[i:])      # keep preposition in highlight
                return before, hl_part, ""

        # Strategy 2: colon split — "Title: Subtitle"
        if ":" in title:
            parts = title.split(":", 1)
            return parts[0].strip(), parts[1].strip(), ""

        # Strategy 3: split after first 2 words
        if len(words) > 3:
            return " ".join(words[:2]), " ".join(words[2:]), ""

        # Strategy 4: whole title as highlight
        return "", title, ""

    before, hl_text, after = _find_split(words)

    first = True
    for part_type, part_text in [("before", before), ("hl", hl_text), ("after", after)]:
        if not part_text.strip():
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = part_text
        run.font.name = h_font
        run.font.bold = True
        if part_type == "hl":
            run.font.size = Pt(hl_size)
            run.font.color.rgb = highlight_color
        elif part_type == "before":
            run.font.size = Pt(base_size)
            run.font.color.rgb = base_color
        else:
            run.font.size = Pt(after_size)
            run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            run.font.bold = False
    return box


def _bullet_block(slide, left, top, w, h, items, size=16, color=None,
                  font="Calibri", spacing=12, max_items=7, bullet_char="\u2022"):
    color = color or RGBColor(0x52, 0x52, 0x5B)
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items[:max_items]):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"{bullet_char}  {_clean(str(item))}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(spacing)
        p.line_spacing = Pt(size + 8)


def _page_num(slide, num, total, color=None):
    color = color or RGBColor(0x71, 0x71, 0x7A)
    _textbox(slide, 12.2, 7.1, 0.8, 0.3, f"{num}/{total}", 8, color, align=PP_ALIGN.RIGHT)


def _info_box(slide, left, top, w, h, header, body_text, header_color, body_color,
              bg_color, font_h="Calibri", font_b="Calibri"):
    """Colored header box with body content — McKinsey info box style."""
    _rect(slide, left, top, w, h, bg_color)
    _rect(slide, left, top, w, 0.38, header_color)
    _textbox(slide, left + 0.12, top + 0.05, w - 0.2, 0.3,
             _clean(str(header))[:45], 12, RGBColor(0xFF, 0xFF, 0xFF), True, font=font_h)
    _textbox(slide, left + 0.12, top + 0.42, w - 0.2, h - 0.5,
             _clean(str(body_text))[:220], 12, body_color, False, font=font_b)


def _bar_chart_rects(slide, left, top, w, h, data_pairs, bar_color, label_color,
                     bg_color=None, font="Calibri"):
    """Draw a minimal bar chart using rectangles. data_pairs = [(label, value), ...]"""
    if not data_pairs:
        return
    max_val = max(max(v for _, v in data_pairs), 1)
    chart_h = h * 0.72
    bar_area_w = w
    n = len(data_pairs)
    slot_w = bar_area_w / n
    bar_w = slot_w * 0.55
    gap_x = slot_w * 0.225

    if bg_color:
        _rect(slide, left, top, w, h, bg_color)

    for i, (label, val) in enumerate(data_pairs):
        bar_h = max((val / max_val) * chart_h, 0.04)
        bx = left + i * slot_w + gap_x
        by = top + h - 0.28 - bar_h
        _rect(slide, bx, by, bar_w, bar_h, bar_color)
        _textbox(slide, bx - gap_x * 0.5, top + h - 0.25, slot_w, 0.22,
                 str(label), 7, label_color, False, PP_ALIGN.CENTER, font)


def _chip_row(slide, left, top, items, chip_colors, font="Calibri"):
    """Render a row of small colored chip/tag elements."""
    x = left
    for i, item in enumerate(items[:6]):
        chip_w = min(len(str(item)) * 0.09 + 0.35, 1.6)
        c = chip_colors[i % len(chip_colors)]
        _rect(slide, x, top, chip_w, 0.32, c)
        _textbox(slide, x + 0.08, top + 0.04, chip_w - 0.12, 0.24,
                 _clean(str(item))[:18], 10, RGBColor(0xFF, 0xFF, 0xFF), True, PP_ALIGN.CENTER, font)
        x += chip_w + 0.12


def _section_cat_label(slide, right_x, top_y, category, color, font):
    """Render the small section category label in top-right area."""
    _textbox(slide, right_x, top_y, 4.0, 0.28, category.upper(), 8,
             color, True, PP_ALIGN.RIGHT, font)


# ── Section category derivation ─────────────────────────────────

_CAT_KEYWORDS = [
    (["neuroscience", "brain", "neural", "cortex", "hippocampus"], "NEUROSCIENCE"),
    (["landmark research", "study", "mit", "stanford", "harvard", "ucl"], "LANDMARK RESEARCH"),
    (["global", "market", "adoption", "scale", "statistic"], "THE SCALE OF ADOPTION"),
    (["cognitive debt", "hidden cost", "debt"], "COGNITIVE DEBT"),
    (["critical thinking", "evidence", "decline", "atrophy"], "CRITICAL THINKING"),
    (["vulnerable", "developing", "risk group", "youth", "age"], "HIGHEST RISK GROUPS"),
    (["trade-off", "efficiency", "vs cognition", "benefit"], "RISK / BENEFIT ANALYSIS"),
    (["recommendation", "strategic", "strategy", "implement"], "RECOMMENDATIONS"),
    (["source", "reference", "citation"], "SOURCES & REFERENCES"),
    (["executive", "summary", "overview"], "EXECUTIVE SUMMARY"),
    (["understanding", "definition", "what is", "introduction"], "UNDERSTANDING THE TECHNOLOGY"),
    (["responsible", "integration", "balance", "neuroscience of"], "RESPONSIBLE AI"),
]

def _derive_section_cat(title: str, layout: str) -> str:
    t = title.lower()
    for keywords, cat in _CAT_KEYWORDS:
        if any(k in t for k in keywords):
            return cat
    layout_defaults = {
        "data_metrics": "KEY METRICS",
        "warning_callout": "CRITICAL FINDINGS",
        "comparison_table": "COMPARATIVE ANALYSIS",
        "benefits_risks": "RISK / BENEFIT ANALYSIS",
        "recommendations": "RECOMMENDATIONS",
        "research_study": "LANDMARK RESEARCH",
        "research_citations": "EVIDENCE & RESEARCH",
        "concept_cards": "UNDERSTANDING THE TECHNOLOGY",
        "anatomy_diagram": "SYSTEM BREAKDOWN",
        "definition_boxes": "UNDERSTANDING THE TECHNOLOGY",
        "two_column_info": "KEY CONCEPTS",
        "executive": "EXECUTIVE SUMMARY",
        "sources": "SOURCES & REFERENCES",
    }
    return layout_defaults.get(layout, "RESEARCH & INSIGHTS")


# ── Main entry ─────────────────────────────────────────────────

def generate_pptx(
    content: str,
    *,
    title: str = "Presentation",
    theme: Literal["light", "dark"] = "light",
    brand: Any = None,
    agent_role: str = "CXO",
    sources: list[str] | None = None,
    add_title_slide: bool = True,
    add_closing_slide: bool = True,
    creative_director: Any = None,
    document_type: str = "presentation",
    subtitle: str = "",
    slide_spec: list[dict[str, Any]] | None = None,
    brand_domain: str = "",
) -> Path:
    DATA_DIR.mkdir(parents=True, exist_ok=True)

    # ── Brand / colour resolution ───────────────────────────────
    if creative_director:
        pri_hex = creative_director.get_primary_color()
        sec_hex = creative_director.get_secondary_color()
        h_font = creative_director.get_heading_font()
        b_font = creative_director.get_body_font()
    else:
        # Premium dark-navy + gold default palette
        pri_hex, sec_hex = "#0F172A", "#F59E0B"
        h_font, b_font = "Calibri", "Calibri"

    if brand:
        if getattr(brand, "primary_color", None):
            pri_hex = brand.primary_color
        if getattr(brand, "secondary_color", None):
            sec_hex = brand.secondary_color
        if getattr(brand, "heading_font", None):
            h_font = brand.heading_font
        if getattr(brand, "body_font", None):
            b_font = brand.body_font

    PRI   = _hex(pri_hex)
    SEC   = _hex(sec_hex)
    GOLD  = RGBColor(0xF5, 0x9E, 0x0B)   # always-gold accent
    DARK  = RGBColor(0x0B, 0x0F, 0x1A)   # deepest navy bg
    DARK2 = RGBColor(0x13, 0x18, 0x28)   # card bg
    DARK3 = RGBColor(0x1E, 0x24, 0x3A)   # lighter card bg
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    OFFWH = RGBColor(0xF8, 0xF9, 0xFA)
    LGRAY = RGBColor(0xE2, 0xE8, 0xF0)
    MGRAY = RGBColor(0x94, 0xA3, 0xB8)
    DGRAY = RGBColor(0x64, 0x74, 0x8B)
    TDARK = RGBColor(0x0F, 0x17, 0x2A)   # text on white
    TBODY = RGBColor(0x33, 0x41, 0x55)   # body text on white
    GREEN = RGBColor(0x10, 0xB9, 0x81)
    RED   = RGBColor(0xEF, 0x44, 0x44)
    ORNG  = RGBColor(0xF5, 0x9E, 0x0B)
    CYAN  = RGBColor(0x06, 0xB6, 0xD4)
    PURP  = RGBColor(0x8B, 0x5C, 0xF6)

    # Brand label
    brand_label = ""
    if brand:
        brand_label = (getattr(brand, "company_name", "") or brand_domain or "").upper()
    if not brand_label and brand_domain:
        brand_label = re.sub(r'^www\.', '', brand_domain).upper().replace(".", "")
    if not brand_label:
        brand_label = "RESEARCH & INSIGHTS"

    # ── Sections setup ──────────────────────────────────────────
    if slide_spec:
        sections = [
            {
                "title":            s.get("section_title", ""),
                "body":             "",
                "bullets":          s.get("bullets", []),
                "layout":           s.get("layout", "content_bullets"),
                "visual_treatment": s.get("visual_treatment", "none"),
                "icon":             s.get("icon", "•"),
                "table_data":       s.get("table_data"),
                "metrics":          s.get("metrics"),
                "benefits":         s.get("benefits"),
                "risks":            s.get("risks"),
                "warning_text":     s.get("warning_text"),
                "quote":            s.get("quote"),
                "quote_attribution":s.get("quote_attribution"),
                "section_category": s.get("section_category", ""),
                "col_headers":      s.get("col_headers"),
                "study_design":     s.get("study_design"),
                "findings":         s.get("findings"),
            }
            for s in slide_spec
            if s.get("section_title") or s.get("bullets")
        ]
        if not sections:
            sections = [{"title": title, "body": "", "bullets": ["Key points"],
                         "layout": "content_bullets", "section_category": ""}]
    else:
        sections = _parse_sections(content)
        for idx, s in enumerate(sections):
            s.update({"layout": "content_bullets", "visual_treatment": "none",
                       "icon": "•", "table_data": None, "metrics": None,
                       "benefits": None, "risks": None, "warning_text": None,
                       "quote": None, "quote_attribution": None,
                       "section_category": "", "col_headers": None,
                       "study_design": None, "findings": None})

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def set_bg(slide, c):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = c

    has_agenda = any("agenda" in (s.get("title") or "").lower() for s in sections)
    auto_agenda = not has_agenda and len(sections) >= 5
    total = (len(sections)
             + (1 if add_title_slide else 0)
             + (1 if auto_agenda else 0)
             + (1 if add_closing_slide else 0))
    sn = 0

    # ── TITLE SLIDE ─────────────────────────────────────────────
    if add_title_slide:
        sn += 1
        s = prs.slides.add_slide(blank)
        set_bg(s, DARK)
        # Top accent bar
        _rect(s, 0, 0, 13.333, 0.07, SEC)
        # Right accent column
        _rect(s, 10.4, 0.07, 2.933, 7.43, DARK2)
        _rect(s, 10.4, 0.07, 0.06, 7.43, GOLD)
        # Bottom bar
        _rect(s, 0, 6.8, 10.4, 0.7, DARK2)
        _rect(s, 0, 6.8, 10.4, 0.05, PRI)
        # Brand badge (gold rectangle with brand label)
        _rect(s, 0.55, 0.28, 1.1, 0.42, GOLD)
        _textbox(s, 0.58, 0.33, 1.04, 0.32, brand_label[:6], 14, DARK, True,
                 PP_ALIGN.CENTER, h_font)
        _textbox(s, 1.85, 0.33, 3.5, 0.32, "RESEARCH & INSIGHTS", 9, GOLD, True,
                 PP_ALIGN.LEFT, h_font)
        # Multi-tone main title
        _rich_title(s, 0.55, 0.9, 9.6, 4.2, title, h_font, GOLD,
                    base_color=WHITE, base_size=32, hl_size=54, after_size=26)
        # Separator line
        _rect(s, 0.55, 5.2, 3.8, 0.05, SEC)
        # Subtitle
        sub = subtitle or "A research-informed analysis of cognitive risks and strategic recommendations."
        _textbox(s, 0.55, 5.35, 9.6, 0.65, sub, 15, MGRAY, font=b_font)
        # Date + brand
        yr = dt.datetime.now().year
        _textbox(s, 0.55, 6.1, 5.5, 0.3,
                 dt.datetime.now().strftime("%B %d, %Y"), 11, DGRAY, font=b_font)
        _textbox(s, 0.55, 6.45, 9.5, 0.28,
                 f"{yr} — {brand_label} Intelligence Series", 11, DGRAY, font=b_font)
        # Decorative dots on right panel
        _rect(s, 10.95, 3.4, 0.22, 0.22, GOLD)
        _rect(s, 11.55, 3.9, 0.18, 0.18, SEC)
        _rect(s, 11.2, 4.5, 0.14, 0.14, MGRAY)
        _page_num(s, sn, total, DGRAY)

    # ── AGENDA ───────────────────────────────────────────────────
    if auto_agenda:
        sn += 1
        s = prs.slides.add_slide(blank)
        set_bg(s, OFFWH)
        _rect(s, 0, 0, 0.08, 7.5, SEC)
        _textbox(s, 0.55, 0.45, 5, 0.75, "Agenda", 32, TDARK, True, font=h_font)
        _rect(s, 0.55, 1.28, 2.2, 0.05, SEC)
        y = 1.55
        for i, sec in enumerate(sections):
            if y > 6.5:
                break
            _textbox(s, 0.55, y, 0.75, 0.42, f"{i+1:02d}", 24, SEC, True, font=h_font)
            _textbox(s, 1.5, y + 0.04, 11, 0.38, _clean(sec.get("title", "")), 17, TBODY, font=b_font)
            y += 0.52
        _page_num(s, sn, total, MGRAY)

    # ── CONTENT SLIDES ───────────────────────────────────────────
    for idx, sec in enumerate(sections):
        sec_title  = sec.get("title") or "Overview"
        bullets    = sec.get("bullets") or []
        stype      = sec.get("layout") or "content_bullets"
        sec_cat    = sec.get("section_category") or _derive_section_cat(sec_title, stype)
        sn += 1
        s = prs.slides.add_slide(blank)

        def hdr_dark(category: str):
            """Standard dark-slide header block."""
            set_bg(s, DARK)
            _rect(s, 0, 0, 13.333, 0.06, SEC)
            _textbox(s, 9.3, 0.14, 3.7, 0.28, brand_label, 8, DGRAY,
                     False, PP_ALIGN.RIGHT, b_font)
            _textbox(s, 0.5, 0.14, 5, 0.28, category, 8, MGRAY, True, font=h_font)
            _textbox(s, 0.5, 0.5, 10.5, 0.72, sec_title, 30, WHITE, True, font=h_font)
            _rect(s, 0.5, 1.28, 2.0, 0.04, GOLD)

        def hdr_light(category: str):
            """Standard light-slide header block."""
            set_bg(s, WHITE)
            _rect(s, 0, 0, 13.333, 0.06, SEC)
            _textbox(s, 9.3, 0.14, 3.7, 0.28, brand_label, 8, MGRAY,
                     False, PP_ALIGN.RIGHT, b_font)
            _textbox(s, 0.5, 0.14, 5, 0.28, category, 8, SEC, True, font=h_font)
            _textbox(s, 0.5, 0.5, 10.5, 0.65, sec_title, 28, TDARK, True, font=h_font)
            _rect(s, 0.5, 1.22, 2.0, 0.04, SEC)

        # ── data_metrics ────────────────────────────────────────
        if stype == "data_metrics":
            metrics = sec.get("metrics") or []
            if not metrics:
                pat = re.compile(r'(-?[\d.]+%|-?\$[\d.]+[BMKTbmkt]?|-?[\d.]+[xX]|[\d,]+[BMKTbmkt+])')
                for b in bullets[:8]:
                    if re.fullmatch(r'\s*20\d\d\s*', b):
                        continue
                    for m in pat.finditer(str(b)):
                        v = m.group(1)
                        if re.fullmatch(r'20\d\d', v.strip('+-')):
                            continue
                        if v.startswith('-') and m.start() > 0 and str(b)[m.start()-1].isdigit():
                            continue
                        lbl = _clean(str(b)).replace(v, "").strip(" :-–")[:55]
                        if lbl and len(lbl) > 3:
                            metrics.append({"value": v, "label": lbl})
                        break
            hdr_dark(sec_cat)
            if metrics:
                count = min(len(metrics), 8)
                if count <= 4:
                    slot_w = 13.333 / count
                    for i, m in enumerate(metrics[:count]):
                        xp = i * slot_w + slot_w * 0.07
                        cw = slot_w * 0.86
                        val_s = str(m.get("value", ""))
                        ac = RED if val_s.startswith("-") else GOLD
                        _rect(s, xp, 1.42, cw, 2.55, DARK3)
                        _rect(s, xp, 1.42, cw, 0.06, ac)
                        _textbox(s, xp, 1.55, cw, 1.15, val_s[:12], 52, ac,
                                 True, PP_ALIGN.CENTER, h_font)
                        _textbox(s, xp, 2.75, cw, 0.9, _clean(str(m.get("label", "")))[:60],
                                 12, MGRAY, False, PP_ALIGN.CENTER, b_font)
                    # Bar chart for adoption/scale slides — extract year/value pairs from bullets
                    if any(k in sec_title.lower() for k in ["scale", "global", "market", "adoption", "growth"]):
                        yr_pat = re.compile(r'\b(20\d\d)\b.*?(\d[\d,.]*)\s*([BMKTbmkt%]|billion|million|trillion)?')
                        chart_pairs: list[tuple[str, float]] = []
                        for bl in bullets:
                            ym = yr_pat.search(bl)
                            if ym:
                                try:
                                    raw_v = ym.group(2).replace(",", "")
                                    mult = {"b": 1e9, "m": 1e6, "k": 1e3, "t": 1e12,
                                            "billion": 1e9, "million": 1e6, "trillion": 1e12}
                                    sfx = (ym.group(3) or "").lower()
                                    v = float(raw_v) * mult.get(sfx, 1)
                                    chart_pairs.append((ym.group(1), v))
                                except (ValueError, AttributeError):
                                    pass
                        # Fallback: use metric values as bar heights if no year pairs found
                        if not chart_pairs and metrics:
                            for mi, m in enumerate(metrics[:6]):
                                try:
                                    v_raw = re.sub(r'[^0-9.]', '', str(m.get("value", "0")))
                                    chart_pairs.append((str(m.get("value", f"#{mi+1}")), float(v_raw) if v_raw else mi + 1))
                                except ValueError:
                                    chart_pairs.append((str(m.get("value", "")), mi + 1.0))
                        # Need ≥2 year pairs for a meaningful timeline chart
                        if len(chart_pairs) < 2 and metrics:
                            chart_pairs = []
                            for mi, m in enumerate(metrics[:6]):
                                try:
                                    v_raw = re.sub(r'[^0-9.]', '', str(m.get("value", "0")))
                                    chart_pairs.append((str(m.get("value", f"#{mi+1}")),
                                                        float(v_raw) if v_raw else float(mi + 1)))
                                except ValueError:
                                    chart_pairs.append((str(m.get("value", "")), float(mi + 1)))
                        if chart_pairs:
                            _bar_chart_rects(s, 0.5, 4.15, 12.3, 2.35, chart_pairs[:8], GOLD, DGRAY, None, b_font)
                            # Dynamic chart label derived from slide title
                            chart_lbl = f"{sec_title} — Source: Industry Reports {dt.datetime.now().year}"
                            _textbox(s, 0.5, 6.6, 12.3, 0.28, chart_lbl, 9, DGRAY,
                                     False, PP_ALIGN.CENTER, b_font)
                    else:
                        ctx = [b for b in bullets if not re.search(r'[\d%$]', b)][:2]
                        if ctx:
                            _textbox(s, 0.5, 4.2, 12.3, 0.45,
                                     "  •  ".join(_clean(c) for c in ctx),
                                     11, DGRAY, False, PP_ALIGN.CENTER, b_font)
                else:
                    # 2-row layout
                    slot_w = 13.333 / 4
                    for row_i, row in enumerate([metrics[:4], metrics[4:8]]):
                        yb = 1.35 + row_i * 2.65
                        for i, m in enumerate(row):
                            xp = i * slot_w + slot_w * 0.05
                            cw = slot_w * 0.9
                            val_s = str(m.get("value", ""))
                            ac = RED if val_s.startswith("-") else GOLD
                            _rect(s, xp, yb, cw, 2.4, DARK3)
                            _rect(s, xp, yb, cw, 0.05, ac)
                            _textbox(s, xp, yb + 0.12, cw, 1.0, val_s[:12],
                                     36, ac, True, PP_ALIGN.CENTER, h_font)
                            _textbox(s, xp, yb + 1.18, cw, 0.95,
                                     _clean(str(m.get("label", "")))[:60],
                                     11, MGRAY, False, PP_ALIGN.CENTER, b_font)
            else:
                _bullet_block(s, 0.5, 1.45, 12, 5.0, bullets[:7], 17, MGRAY, b_font, 14, 7, "\u25B8")
            _page_num(s, sn, total, DGRAY)

        # ── research_study ──────────────────────────────────────
        elif stype == "research_study":
            hdr_dark(sec_cat)
            study_design = sec.get("study_design") or []
            findings     = sec.get("findings") or []
            quote_text   = sec.get("quote") or ""

            # If no structured data, extract from bullets
            if not study_design:
                study_design = [b for b in bullets if b.lower().startswith("study design")]
                if not study_design:
                    study_design = bullets[:3]
            if not findings:
                findings = [b for b in bullets if re.search(r'\d+%', b)][:3]
                if not findings:
                    findings = bullets[len(study_design):len(study_design)+3]
            if not quote_text:
                qb = [b for b in bullets if b.lower().startswith("quote") or
                      ('"' in b and len(b) > 60)]
                quote_text = qb[0].replace("Quote from lead researcher:", "").strip().strip('"') if qb else ""

            # Study citation banner
            _rect(s, 0.4, 1.45, 12.5, 0.55, DARK3)
            _rect(s, 0.4, 1.45, 0.07, 0.55, GOLD)
            study_banner = (study_design[0].replace("Study Design:", "").strip()[:120]
                            if study_design else sec_title)
            _textbox(s, 0.65, 1.52, 12.0, 0.38, study_banner, 12, MGRAY, False, font=b_font)

            # Left: study design panel
            _rect(s, 0.4, 2.1, 4.6, 4.1, DARK3)
            _textbox(s, 0.55, 2.18, 4.2, 0.35, "STUDY DESIGN", 10, GOLD, True, font=h_font)
            sd_bullets = (study_design[0].replace("Study Design:", "").strip().split(",")
                          if study_design else bullets[:4])
            ysd = 2.6
            for sd_b in sd_bullets[:5]:
                _textbox(s, 0.55, ysd, 4.2, 0.48, f"\u25B8  {_clean(sd_b.strip())[:60]}",
                         12, MGRAY, False, font=b_font)
                ysd += 0.52

            # Right: 3 finding columns — derive column titles from findings content
            def _short_col_title(finding_text: str) -> str:
                """Extract a 2-4 word column title from a finding bullet."""
                t = _clean(str(finding_text))
                # Try colon-prefixed label first (e.g. "Critical thinking assessment scores: ...")
                if ":" in t:
                    lbl = t.split(":")[0].strip()
                    words = lbl.split()
                    if 1 <= len(words) <= 5:
                        return " ".join(words[:4]).title()
                # Otherwise take first 3 meaningful words
                words = [w for w in t.split() if len(w) > 2 and not w[0].isdigit()]
                return " ".join(words[:3]).title() if words else "Key Finding"
            col_titles = [_short_col_title(findings[ci]) if ci < len(findings) else f"Finding {ci+1}"
                          for ci in range(3)]
            col_colors = [RED, ORNG, CYAN]
            cx_positions = [5.2, 8.05, 10.9]
            for ci in range(min(3, len(findings))):
                cx = cx_positions[ci]
                ftext = _clean(str(findings[ci]))
                val_m = re.search(r'(-?\d+\.?\d*%|-?\d+\.?\d*[xX])', ftext)
                val_disp = val_m.group(1) if val_m else ""
                col_lbl = col_titles[ci] if ci < len(col_titles) else f"Finding {ci+1}"
                cc = col_colors[ci]
                _rect(s, cx, 2.1, 2.9, 0.4, DARK3)
                _textbox(s, cx + 0.1, 2.15, 2.7, 0.3, col_lbl, 10, cc, True, font=h_font)
                if val_disp:
                    _textbox(s, cx, 2.58, 2.9, 0.95, val_disp, 48, cc,
                              True, PP_ALIGN.CENTER, h_font)
                _textbox(s, cx + 0.1, 3.6, 2.7, 2.55, ftext[:150], 12, MGRAY, False, font=b_font)

            # Bottom: quote bar
            if quote_text:
                _rect(s, 0.4, 6.3, 12.5, 0.88, RGBColor(0x14, 0x18, 0x2E))
                _textbox(s, 0.55, 6.35, 0.5, 0.55, "\u201C", 28, GOLD, True, font=h_font)
                _textbox(s, 1.15, 6.42, 11.0, 0.68,
                         _clean(str(quote_text))[:220], 12, WHITE, False, font=b_font)
            _page_num(s, sn, total, DGRAY)

        # ── concept_cards ────────────────────────────────────────
        # 3-card layout: definition box + concept cards + examples footer
        elif stype == "concept_cards":
            hdr_dark(sec_cat)
            definition = _clean(str(sec.get("definition") or ""))
            concepts   = sec.get("concepts") or []
            footer     = _clean(str(sec.get("footer") or ""))

            # Definition banner
            def_bg = RGBColor(0x0D, 0x1F, 0x1A)  # very dark green
            _rect(s, 0.4, 1.42, 12.5, 0.82, def_bg)
            _rect(s, 0.4, 1.42, 0.07, 0.82, GREEN)
            _textbox(s, 0.65, 1.47, 1.6, 0.28, "DEFINITION", 9, GREEN, True, font=h_font)
            if definition:
                _textbox(s, 0.65, 1.72, 12.0, 0.44, definition[:280], 13, MGRAY, False, font=b_font)

            # 3 concept cards
            card_colors = [SEC, CYAN, GOLD]
            card_w, gap = 4.05, 0.13
            for ci, concept in enumerate(concepts[:3]):
                cx = 0.4 + ci * (card_w + gap)
                cc = card_colors[ci % 3]
                name = _clean(str(concept.get("name", f"Concept {ci+1}")))[:30]
                desc = _clean(str(concept.get("description", "")))[:240]
                exs  = _clean(str(concept.get("examples", "")))[:80]
                _rect(s, cx, 2.36, card_w, 3.5, DARK3)
                _rect(s, cx, 2.36, card_w, 0.06, cc)
                _textbox(s, cx + 0.18, 2.48, card_w - 0.3, 0.55, name, 20, cc, True, font=h_font)
                if desc:
                    _textbox(s, cx + 0.18, 3.1, card_w - 0.3, 1.9, desc, 12, MGRAY, False, font=b_font)
                if exs:
                    _rect(s, cx, 5.62, card_w, 0.24, RGBColor(0x0B, 0x10, 0x20))
                    _textbox(s, cx + 0.1, 5.65, card_w - 0.15, 0.2,
                             exs, 8, DGRAY, False, font=b_font)

            # Footer examples row
            if footer:
                _textbox(s, 0.5, 6.05, 12.3, 0.38, footer, 11, DGRAY, False, font=b_font)
            _page_num(s, sn, total, DGRAY)

        # ── anatomy_diagram ──────────────────────────────────────
        # Component breakdown: left = labeled list, right = 2 effect panels
        elif stype == "anatomy_diagram":
            hdr_light(sec_cat)
            components   = sec.get("components") or []
            right_panels = sec.get("right_panels") or []

            comp_colors = [PRI, GREEN, PURP, ORNG, CYAN]

            # Left: labeled component list
            y = 1.38
            for ci, comp in enumerate(components[:5]):
                cc = comp_colors[ci % len(comp_colors)]
                name  = _clean(str(comp.get("name", "")))[:40]
                funcs = comp.get("functions") or []
                row_h = 0.42 + len(funcs[:2]) * 0.28
                _rect(s, 0.4, y, 0.08, row_h, cc)
                _textbox(s, 0.62, y + 0.02, 5.1, 0.38, name, 14, TDARK, True, font=h_font)
                for fi, func in enumerate(funcs[:2]):
                    _textbox(s, 0.62, y + 0.40 + fi * 0.28, 5.1, 0.26,
                             f"  \u25B8  {_clean(str(func))[:70]}", 11, TBODY, False, font=b_font)
                y += row_h + 0.28

            # Vertical divider
            _rect(s, 6.05, 1.38, 0.03, 5.4, LGRAY)

            # Right: 2 info panels
            panel_colors = [ORNG, RED]
            py = 1.38
            for pi, panel in enumerate(right_panels[:2]):
                ph = 2.55
                pc = panel_colors[pi % len(panel_colors)]
                hdr_txt  = _clean(str(panel.get("header", f"Finding {pi+1}")))[:50]
                body_txt = _clean(str(panel.get("body", "")))[:300]
                _info_box(s, 6.25, py, 6.85, ph, hdr_txt, body_txt, pc, TBODY, OFFWH, h_font, b_font)
                py += ph + 0.28
            _page_num(s, sn, total, MGRAY)

        # ── research_citations ───────────────────────────────────
        # Evidence slide: study list + negative metrics + domain chips
        elif stype == "research_citations":
            hdr_light(sec_cat)
            studies      = sec.get("studies") or []
            metrics_list = sec.get("metrics") or []
            domains      = sec.get("domains") or []
            footer_quote = _clean(str(sec.get("footer_quote") or ""))

            # Left: evidence base panel (0.4 to 4.1)
            _rect(s, 0.4, 1.38, 3.55, 5.35, RGBColor(0xF8, 0xF9, 0xFA))
            _rect(s, 0.4, 1.38, 3.55, 0.38, PRI)
            _textbox(s, 0.55, 1.44, 3.2, 0.28, "EVIDENCE BASE", 10, WHITE, True, font=h_font)
            sy = 1.85
            for si, study in enumerate(studies[:3]):
                name   = _clean(str(study.get("name", "")))[:55]
                detail = _clean(str(study.get("detail", "")))[:160]
                _rect(s, 0.4, sy, 0.05, 1.05, SEC)
                _textbox(s, 0.6, sy + 0.04, 3.2, 0.34, name, 13, TDARK, True, font=h_font)
                if detail:
                    _textbox(s, 0.6, sy + 0.38, 3.2, 0.6, detail, 10, TBODY, False, font=b_font)
                sy += 1.2

            # Center: stacked metric boxes (4.25 to 7.0)
            mx = 4.25
            mw = 2.55
            my = 1.38
            nm = min(len(metrics_list), 3)
            mh = 5.35 / max(nm, 1) - 0.12
            for mi, m in enumerate(metrics_list[:3]):
                val   = str(m.get("value", ""))
                label = _clean(str(m.get("label", "")))[:50]
                mc    = RED if val.startswith("-") else SEC
                _rect(s, mx, my, mw, mh, DARK3)
                _rect(s, mx, my, mw, 0.06, mc)
                _textbox(s, mx, my + 0.12, mw, mh * 0.55, val[:10],
                         int(min(52, 4.5 * mw)), mc, True, PP_ALIGN.CENTER, h_font)
                _textbox(s, mx + 0.1, my + mh * 0.6, mw - 0.2, mh * 0.38,
                         label, 10, MGRAY, False, PP_ALIGN.CENTER, b_font)
                my += mh + 0.12

            # Right: key finding text (7.1 to 12.93)
            right_bullets = [b for b in bullets if not re.search(r'\(\d{4}\)', b)][:4]
            _rect(s, 7.1, 1.38, 5.7, 5.35, OFFWH)
            _rect(s, 7.1, 1.38, 5.7, 0.38, SEC)
            _textbox(s, 7.25, 1.43, 5.3, 0.28, "KEY FINDINGS", 10, WHITE, True, font=h_font)
            rb_y = 1.9
            for rb in right_bullets:
                _rect(s, 7.1, rb_y, 0.05, 0.52, GOLD)
                _textbox(s, 7.28, rb_y + 0.04, 5.35, 0.52,
                         _clean(str(rb))[:140], 12, TBODY, False, font=b_font)
                rb_y += 0.65

            # Domain chip row
            if domains:
                _chip_row(s, 0.4, 6.88, domains[:6],
                          [PRI, SEC, GREEN, ORNG, PURP, CYAN], b_font)

            # Footer quote
            if footer_quote:
                _rect(s, 0.4, 6.5, 12.5, 0.32, DARK3)
                _textbox(s, 0.55, 6.55, 12.1, 0.24,
                         f"\u201C {footer_quote[:180]} \u201D", 9, DGRAY, False, font=b_font)
            _page_num(s, sn, total, MGRAY)

        # ── definition_boxes ────────────────────────────────────
        elif stype == "definition_boxes":
            hdr_dark(sec_cat)
            colon_bullets = [b for b in bullets if ":" in b]
            # Optional definition box (first colon bullet treated as key term def)
            if colon_bullets:
                parts0 = colon_bullets[0].split(":", 1)
                _rect(s, 0.4, 1.45, 12.5, 0.88, RGBColor(0x18, 0x1A, 0x2A))
                _rect(s, 0.4, 1.45, 0.07, 0.88, GOLD)
                _textbox(s, 0.65, 1.5, 2.5, 0.35,
                         _clean(parts0[0])[:30], 13, GOLD, True, font=h_font)
                _textbox(s, 3.25, 1.5, 9.5, 0.7,
                         _clean(parts0[1])[:200], 13, MGRAY, False, font=b_font)

            # 3 concept boxes
            boxes = colon_bullets[1:4] if len(colon_bullets) >= 4 else colon_bullets[:3]
            if not boxes:
                boxes = bullets[:3]
            box_colors = [SEC, CYAN, GOLD]
            bw = 12.4 / max(len(boxes), 1)
            for bi, cb in enumerate(boxes[:3]):
                parts = cb.split(":", 1) if ":" in cb else [cb, ""]
                bx = 0.45 + bi * (bw + 0.04)
                bc = box_colors[bi % 3]
                _rect(s, bx, 2.45, bw, 3.0, DARK3)
                _rect(s, bx, 2.45, bw, 0.07, bc)
                _textbox(s, bx + 0.18, 2.58, bw - 0.3, 0.48,
                         _clean(parts[0])[:30], 18, bc, True, font=h_font)
                _textbox(s, bx + 0.18, 3.1, bw - 0.3, 1.75,
                         _clean(parts[1].strip())[:200], 13, MGRAY, False, font=b_font)

            # Examples row at bottom
            non_colon = [b for b in bullets if ":" not in b or len(b.split(":", 1)[0]) > 30]
            if non_colon:
                _textbox(s, 0.5, 5.65, 12.3, 0.35,
                         "Examples: " + "  ·  ".join(_clean(b)[:30] for b in non_colon[:5]),
                         11, DGRAY, False, font=b_font)
            _page_num(s, sn, total, DGRAY)

        # ── two_column_info ─────────────────────────────────────
        elif stype == "two_column_info":
            hdr_light(sec_cat)
            colon_bullets = [b for b in bullets if ":" in b and len(b.split(":", 1)[0]) < 35]
            accent_colors = [PRI, SEC, ORNG, GREEN, PURP, CYAN]

            # Left column: labeled list
            y = 1.4
            for bi, cb in enumerate(colon_bullets[:5]):
                parts = cb.split(":", 1)
                ac = accent_colors[bi % len(accent_colors)]
                lbl = _clean(parts[0])[:28]
                desc = _clean(parts[1].strip())[:130] if len(parts) > 1 else ""
                num_m = re.search(r'(\d+[%$xX+]|-\d+%)', cb)
                nv = num_m.group(1) if num_m else ""
                _rect(s, 0.4, y, 0.06, 0.78, ac)
                _textbox(s, 0.6, y, 3.4, 0.36, lbl, 14, TDARK, True, font=h_font)
                if desc:
                    _textbox(s, 0.6, y + 0.35, 5.5, 0.42, desc[:110], 11, TBODY, False, font=b_font)
                if nv:
                    nc = RED if nv.startswith('-') else SEC
                    _textbox(s, 5.4, y, 1.5, 0.78, nv, 28, nc, True, PP_ALIGN.CENTER, h_font)
                y += 0.9

            # Right column: 2 info boxes
            info_bullets = [b for b in bullets if b not in colon_bullets]
            if not info_bullets:
                info_bullets = colon_bullets[5:]
            if len(info_bullets) >= 1:
                _info_box(s, 7.15, 1.4, 5.8, 2.35,
                          info_bullets[0].split(":")[0] if ":" in info_bullets[0] else "Key Finding",
                          info_bullets[0], ORNG, TBODY, OFFWH, h_font, b_font)
            if len(info_bullets) >= 2:
                _info_box(s, 7.15, 3.88, 5.8, 2.35,
                          info_bullets[1].split(":")[0] if ":" in info_bullets[1] else "Neural Impact",
                          info_bullets[1], RED, TBODY, OFFWH, h_font, b_font)
            _page_num(s, sn, total, MGRAY)

        # ── warning_callout ─────────────────────────────────────
        elif stype == "warning_callout":
            warn = sec.get("warning_text") or (bullets[0] if bullets else "Important notice.")
            support = [b for b in bullets if b != warn]
            age_bullets = [b for b in bullets if re.match(r'^(Ages?\s+\d|\d{1,2}[-–]\d{1,2})', b)]

            if len(age_bullets) >= 3:
                # Age-group rows variant (Most Vulnerable)
                set_bg(s, DARK)
                _rect(s, 0, 0, 13.333, 0.07, RED)
                _textbox(s, 9.3, 0.14, 3.7, 0.28, brand_label, 8, DGRAY, False, PP_ALIGN.RIGHT, b_font)
                _textbox(s, 0.5, 0.14, 5, 0.28, sec_cat, 8, RED, True, font=h_font)
                _textbox(s, 0.5, 0.5, 10.5, 0.7, sec_title, 28, WHITE, True, font=h_font)
                _rect(s, 0.5, 1.28, 2.0, 0.04, RED)
                # Warning banner
                _rect(s, 0.4, 1.45, 12.5, 0.55, RGBColor(0x7F, 0x1D, 0x1D))
                _textbox(s, 0.55, 1.5, 0.5, 0.4, "⚠", 18, ORNG, True, font=h_font)
                _textbox(s, 1.1, 1.55, 11.6, 0.38,
                         _clean(str(warn))[:170], 13, RGBColor(0xFE, 0xCA, 0xCA), False, font=b_font)
                # Age rows
                row_bgs = [DARK3, RGBColor(0x1F, 0x14, 0x14), DARK3, RGBColor(0x14, 0x18, 0x2A)]
                lbl_bgs = [RGBColor(0x2D, 0x1B, 0x69), RGBColor(0x44, 0x12, 0x1A),
                           RGBColor(0x1A, 0x32, 0x5C), RGBColor(0x27, 0x2A, 0x1A)]
                y = 2.12
                for ri, ab in enumerate(age_bullets[:4]):
                    parts = ab.split(":", 1)
                    age_lbl = _clean(parts[0])[:18]
                    age_desc = _clean(parts[1].strip())[:190] if len(parts) > 1 else _clean(ab)[:190]
                    _rect(s, 0.4, y, 12.5, 1.2, row_bgs[ri % 4])
                    _rect(s, 0.4, y, 2.1, 1.2, lbl_bgs[ri % 4])
                    _textbox(s, 0.45, y + 0.12, 2.0, 0.95, age_lbl,
                              22, WHITE, True, PP_ALIGN.CENTER, h_font)
                    _textbox(s, 2.7, y + 0.18, 9.9, 0.85, age_desc,
                              13, MGRAY, False, font=b_font)
                    y += 1.26
                non_age = [b for b in bullets if b not in age_bullets and b != warn]
                if non_age:
                    _textbox(s, 0.5, y + 0.08, 12.3, 0.42, _clean(non_age[0])[:200],
                             11, DGRAY, False, font=b_font)
            else:
                # Standard warning
                hdr_light(sec_cat)
                _rect(s, 0.4, 1.38, 12.5, 1.8, RGBColor(0xFF, 0xF3, 0xE0))
                _rect(s, 0.4, 1.38, 0.08, 1.8, ORNG)
                _textbox(s, 0.65, 1.48, 0.95, 0.5, "⚠", 28, ORNG, True, font=h_font)
                _textbox(s, 1.62, 1.48, 11.0, 0.38, "KEY FINDING", 11, ORNG, True, font=h_font)
                _textbox(s, 1.62, 1.88, 11.0, 1.15,
                         _clean(str(warn))[:350], 14, TDARK, False, font=b_font)
                # Support boxes
                if support:
                    n = min(len(support[:4]), 4)
                    bw = 12.2 / n
                    box_cs = [PRI, SEC, ORNG, GREEN]
                    for bi, b in enumerate(support[:4]):
                        bx = 0.4 + bi * bw
                        _rect(s, bx, 3.32, bw - 0.1, 2.0, box_cs[bi % 4])
                        _textbox(s, bx + 0.14, 3.42, bw - 0.25, 1.8,
                                 _clean(str(b))[:140], 13, WHITE, False, font=b_font)
            _page_num(s, sn, total, MGRAY)

        # ── benefits_risks ──────────────────────────────────────
        elif stype == "benefits_risks":
            benefits = sec.get("benefits") or []
            risks    = sec.get("risks") or []
            if not benefits and not risks:
                mid = max(len(bullets) // 2, 1)
                benefits = bullets[:mid]
                risks    = bullets[mid:] if mid < len(bullets) else bullets[:1]
            hdr_light(sec_cat)
            # Column header bars
            _rect(s, 0.4, 1.38, 6.0, 0.5, GREEN)
            _textbox(s, 0.6, 1.42, 5.7, 0.38, "✓  POTENTIAL BENEFITS", 13, WHITE, True, font=h_font)
            _rect(s, 6.9, 1.38, 6.0, 0.5, RED)
            _textbox(s, 7.1, 1.42, 5.7, 0.38, "✗  COGNITIVE RISKS", 13, WHITE, True, font=h_font)
            # Benefits
            yb = 1.98
            for b in benefits[:5]:
                _rect(s, 0.4, yb, 0.06, 0.48, GREEN)
                _textbox(s, 0.62, yb, 5.8, 0.52, _clean(str(b))[:140], 13, TBODY, False, font=b_font)
                yb += 0.58
            # Risks
            yr = 1.98
            for r in risks[:5]:
                _rect(s, 6.9, yr, 0.06, 0.48, RED)
                _textbox(s, 7.12, yr, 5.8, 0.52, _clean(str(r))[:140], 13, TBODY, False, font=b_font)
                yr += 0.58
            # Assessment footer — derive from content, not hardcoded
            assessment_note = ""
            if benefits and risks:
                b_count = len(benefits)
                r_count = len(risks)
                if r_count > b_count:
                    assessment_note = f"{brand_label} Assessment: Risks outweigh benefits — proceed with a structured mitigation strategy."
                elif b_count > r_count:
                    assessment_note = f"{brand_label} Assessment: Benefits predominate — prioritize deployment with appropriate safeguards."
                else:
                    assessment_note = f"{brand_label} Assessment: Balanced trade-off — context-specific evaluation required before adoption."
            else:
                assessment_note = f"{brand_label} Assessment: Weigh these trade-offs carefully against your organization's specific context."
            _rect(s, 0.4, 6.4, 12.5, 0.42, DARK3)
            _textbox(s, 0.6, 6.46, 12.1, 0.32, assessment_note, 11, MGRAY, False, font=b_font)
            _page_num(s, sn, total, MGRAY)

        # ── comparison_table ────────────────────────────────────
        elif stype == "comparison_table":
            tbl = sec.get("table_data") or {}
            headers = tbl.get("headers") or sec.get("col_headers") or ["Before AI", "With Agentic AI"]
            rows    = tbl.get("rows", [])
            if not rows and len(bullets) >= 2:
                mid   = len(bullets) // 2
                left  = bullets[:mid]
                right = bullets[mid:]
                rows  = [[l, r] for l, r in zip(left, right)]
            hdr_light(sec_cat)
            # Column headers
            h0c = GREEN if any(k in str(headers[0]).lower() for k in ["benefit", "positive", "pro"]) else PRI
            h1c = RED   if any(k in str(headers[1]).lower() for k in ["risk", "cost", "con", "negative"]) else SEC
            _rect(s, 0.4, 1.38, 6.1, 0.5, h0c)
            _textbox(s, 0.6, 1.43, 5.8, 0.38, str(headers[0])[:40], 13, WHITE, True, font=h_font)
            _rect(s, 6.9, 1.38, 6.1, 0.5, h1c)
            _textbox(s, 7.1, 1.43, 5.8, 0.38, str(headers[1])[:40], 13, WHITE, True, font=h_font)
            # Rows
            y = 1.95
            for ri, row in enumerate(rows[:6]):
                if len(row) >= 2:
                    rbg = RGBColor(0xF8, 0xF9, 0xFA) if ri % 2 == 0 else WHITE
                    _rect(s, 0.4, y, 12.6, 0.55, rbg)
                    _rect(s, 6.87, y, 0.03, 0.55, LGRAY)
                    _textbox(s, 0.6, y + 0.06, 6.0, 0.45,
                             _clean(str(row[0]))[:100], 13, TBODY, False, font=b_font)
                    _textbox(s, 7.1, y + 0.06, 6.0, 0.45,
                             _clean(str(row[1]))[:100], 13, TBODY, False, font=b_font)
                y += 0.57
            _page_num(s, sn, total, MGRAY)

        # ── recommendations ─────────────────────────────────────
        elif stype == "recommendations":
            hdr_light(sec_cat)
            _rect(s, 0, 6.9, 13.333, 0.6, DARK)
            items = bullets[:4]
            nc_list = [PRI, SEC, GREEN, ORNG]
            cat_labels = ["INDIVIDUALLY", "ORGANIZATIONS", "POLICY", "MEASUREMENT"]
            if len(items) >= 3:
                # 2×2 grid
                grid = [(0.4, 1.38), (6.9, 1.38), (0.4, 4.15), (6.9, 4.15)]
                for gi, (gx, gy) in enumerate(grid[:len(items)]):
                    nc = nc_list[gi % len(nc_list)]
                    btext = _clean(items[gi])
                    parts = btext.split(":", 1)
                    heading = parts[0].strip()[:55]
                    detail  = parts[1].strip()[:240] if len(parts) > 1 else btext[:240]
                    _textbox(s, gx, gy, 6.0, 0.28,
                             cat_labels[gi] if gi < len(cat_labels) else f"0{gi+1}",
                             8, MGRAY, True, font=h_font)
                    _rect(s, gx, gy + 0.3, 0.52, 0.52, nc)
                    _textbox(s, gx + 0.04, gy + 0.34, 0.44, 0.44,
                             str(gi + 1), 21, WHITE, True, PP_ALIGN.CENTER, h_font)
                    _textbox(s, gx + 0.68, gy + 0.33, 5.7, 0.5,
                             heading, 15, TDARK, True, font=h_font)
                    if detail:
                        _textbox(s, gx + 0.68, gy + 0.86, 5.7, 1.85,
                                 detail, 12, TBODY, False, font=b_font)
                    _rect(s, gx, gy + 2.65, 6.1, 0.03, LGRAY)
            else:
                y = 1.42
                for i, b in enumerate(bullets[:5]):
                    nc = nc_list[i % len(nc_list)]
                    parts = _clean(b).split(":", 1)
                    _rect(s, 0.4, y, 0.52, 0.52, nc)
                    _textbox(s, 0.44, y + 0.06, 0.44, 0.4,
                             str(i + 1), 21, WHITE, True, PP_ALIGN.CENTER, h_font)
                    _textbox(s, 1.1, y + 0.06, 11.8, 0.42,
                             parts[0].strip(), 16, TDARK, True, font=h_font)
                    if len(parts) > 1:
                        _textbox(s, 1.1, y + 0.52, 11.8, 0.4,
                                 parts[1].strip()[:140], 13, TBODY, False, font=b_font)
                    y += 1.06
            _page_num(s, sn, total, MGRAY)

        # ── executive ───────────────────────────────────────────
        elif stype == "executive":
            hdr_dark(sec_cat)
            colon_b = [b for b in bullets if ":" in b]
            if len(colon_b) >= 3:
                bw = 12.4 / 3
                box_cs = [SEC, CYAN, GOLD]
                for bi, cb in enumerate(colon_b[:3]):
                    parts = cb.split(":", 1)
                    bx = 0.45 + bi * (bw + 0.04)
                    bc = box_cs[bi % 3]
                    _rect(s, bx, 1.45, bw, 2.8, DARK3)
                    _rect(s, bx, 1.45, bw, 0.06, bc)
                    _textbox(s, bx + 0.16, 1.58, bw - 0.28, 0.48,
                             _clean(parts[0])[:32], 16, bc, True, font=h_font)
                    _textbox(s, bx + 0.16, 2.1, bw - 0.28, 1.9,
                             _clean(parts[1].strip())[:240], 13, MGRAY, False, font=b_font)
                rest = [b for b in bullets if b not in colon_b]
                if rest:
                    _textbox(s, 0.5, 4.45, 12.3, 0.38,
                             "  •  ".join(_clean(r)[:35] for r in rest[:4]),
                             12, DGRAY, False, PP_ALIGN.CENTER, b_font)
            else:
                _bullet_block(s, 0.5, 1.45, 12.0, 5.4, bullets, 17, MGRAY, b_font, 14, 7, "\u25B8")
            _page_num(s, sn, total, DGRAY)

        # ── agenda ──────────────────────────────────────────────
        elif stype == "agenda":
            set_bg(s, OFFWH)
            _rect(s, 0, 0, 0.08, 7.5, SEC)
            _textbox(s, 0.55, 0.45, 5, 0.75, sec_title, 30, TDARK, True, font=h_font)
            _rect(s, 0.55, 1.28, 2.2, 0.05, SEC)
            if bullets:
                y = 1.55
                for i, b in enumerate(bullets[:12]):
                    _textbox(s, 0.55, y, 0.75, 0.42, f"{i+1:02d}", 22, SEC, True, font=h_font)
                    _textbox(s, 1.5, y + 0.04, 11, 0.38, _clean(b), 17, TBODY, font=b_font)
                    y += 0.48
            _page_num(s, sn, total, MGRAY)

        # ── sources ─────────────────────────────────────────────
        elif stype == "sources":
            set_bg(s, OFFWH)
            _rect(s, 0, 0, 13.333, 0.06, SEC)
            _textbox(s, 9.3, 0.14, 3.7, 0.28, brand_label, 8, MGRAY, False, PP_ALIGN.RIGHT, b_font)
            _textbox(s, 0.5, 0.14, 5, 0.28, sec_cat, 8, SEC, True, font=h_font)
            _textbox(s, 0.5, 0.5, 10.5, 0.65, sec_title, 28, TDARK, True, font=h_font)
            _rect(s, 0.5, 1.22, 2.0, 0.04, SEC)
            y = 1.45
            for i, b in enumerate(bullets[:10]):
                _rect(s, 0.4, y, 0.06, 0.38, SEC if i % 2 == 0 else GOLD)
                _textbox(s, 0.65, y, 12.0, 0.42, _clean(b), 13, TBODY, False, font=b_font)
                y += 0.48
            _page_num(s, sn, total, MGRAY)

        # ── quote ───────────────────────────────────────────────
        elif stype == "quote":
            hdr_dark(sec_cat)
            _textbox(s, 0.8, 1.5, 1.5, 1.5, "\u201C", 96, GOLD, True, font=h_font)
            quote = bullets[0] if bullets else ""
            _textbox(s, 2.2, 2.2, 9.5, 2.5, _clean(quote), 24, WHITE, False, font=b_font,
                     line_spacing=36)
            if len(bullets) > 1:
                _textbox(s, 2.2, 5.0, 9.5, 0.5,
                         f"\u2014 {_clean(bullets[1])}", 14, MGRAY, True, font=b_font)
            _page_num(s, sn, total, DGRAY)

        # ── bottom_line ─────────────────────────────────────────
        elif stype == "bottom_line":
            set_bg(s, DARK)
            _rect(s, 0, 0, 13.333, 0.06, SEC)
            _textbox(s, 9.3, 0.14, 3.7, 0.28, brand_label, 8, DGRAY, False, PP_ALIGN.RIGHT, b_font)
            # Large decorative circle
            circle = s.shapes.add_shape(9, Inches(8.2), Inches(0.8), Inches(5.5), Inches(5.5))
            circle.fill.solid()
            circle.fill.fore_color.rgb = DARK2
            circle.line.fill.background()
            # Gold title
            _textbox(s, 0.6, 0.5, 7.5, 0.75, "The Bottom Line", 42, GOLD, True, font=h_font)
            if bullets:
                _textbox(s, 0.6, 1.45, 7.5, 0.55, _clean(bullets[0])[:90], 22, WHITE, False, font=b_font)
                if len(bullets) > 1:
                    _textbox(s, 0.6, 2.05, 7.5, 0.48,
                             _clean(bullets[1])[:90], 18, MGRAY, False, font=b_font)
            # 3 action rows — derive short verb labels from bullet content
            action_bullets = bullets[2:5] if len(bullets) > 2 else bullets[:3]
            _action_verbs = ["ACT", "CONSIDER", "PRIORITIZE", "FOCUS", "IMPLEMENT",
                             "EVALUATE", "MEASURE", "PROTECT", "INVEST", "BUILD"]
            _default_labels = ["THINK", "PROTECT", "INVEST"]
            def _row_label(bullet_text: str, fallback: str) -> str:
                t = _clean(str(bullet_text)).upper()
                for verb in _action_verbs:
                    if verb in t:
                        return verb
                # Use first colon-prefixed keyword if present
                if ":" in bullet_text:
                    lbl = bullet_text.split(":")[0].strip().split()
                    if lbl and len(lbl[0]) <= 10:
                        return lbl[0].upper()
                return fallback
            row_labels = [_row_label(action_bullets[ri], _default_labels[ri])
                          if ri < len(action_bullets) else _default_labels[ri]
                          for ri in range(3)]
            row_colors = [PRI, GREEN, ORNG]
            for ri, (ab, rl, rc) in enumerate(zip(action_bullets, row_labels, row_colors)):
                ry = 3.0 + ri * 1.05
                _rect(s, 0.5, ry, 1.15, 0.78, rc)
                _textbox(s, 0.52, ry + 0.14, 1.11, 0.5,
                         rl, 14, WHITE, True, PP_ALIGN.CENTER, h_font)
                _rect(s, 1.7, ry, 10.5, 0.78, DARK3)
                _textbox(s, 1.88, ry + 0.1, 10.2, 0.6,
                         _clean(str(ab))[:130], 14, MGRAY, False, font=b_font)
            _page_num(s, sn, total, DGRAY)

        # ── two_column (fallback for 7+ bullets) ────────────────
        elif stype == "two_column":
            hdr_light(sec_cat)
            mid = len(bullets) // 2
            _bullet_block(s, 0.5, 1.42, 5.9, 5.4, bullets[:mid], 15, TBODY, b_font, 10, 8)
            _rect(s, 6.65, 1.42, 0.03, 5.0, LGRAY)
            _bullet_block(s, 6.9, 1.42, 5.9, 5.4, bullets[mid:], 15, TBODY, b_font, 10, 8)
            _page_num(s, sn, total, MGRAY)

        # ── content_bullets (default) ───────────────────────────
        else:
            colon_b = [b for b in bullets if ":" in b and len(b.split(":", 1)[0]) < 35]

            if len(colon_b) >= 3:
                # Labeled rows with colored accent chips
                hdr_light(sec_cat)
                ac_colors = [PRI, SEC, ORNG, GREEN, PURP, CYAN]
                y = 1.42
                for bi, cb in enumerate(colon_b[:6]):
                    parts = cb.split(":", 1)
                    lbl  = _clean(parts[0])[:30]
                    desc = _clean(parts[1].strip())[:140] if len(parts) > 1 else ""
                    ac   = ac_colors[bi % len(ac_colors)]
                    num_m = re.search(r'(\d+[%$xX+]|-\d+%)', cb)
                    nv = num_m.group(1) if num_m else ""
                    _rect(s, 0.4, y, 0.06, 0.82, ac)
                    _textbox(s, 0.6, y, 3.5, 0.38, lbl, 14, TDARK, True, font=h_font)
                    if desc:
                        _textbox(s, 0.6, y + 0.36, 5.8, 0.42, desc[:130], 12, TBODY, False, font=b_font)
                    if nv:
                        nc = RED if nv.startswith('-') else SEC
                        _textbox(s, 10.8, y, 2.2, 0.82, nv, 28, nc, True, PP_ALIGN.CENTER, h_font)
                    y += 0.92
                _page_num(s, sn, total, MGRAY)

            elif idx % 3 == 0:
                # Dark variant for visual rhythm
                hdr_dark(sec_cat)
                _bullet_block(s, 0.5, 1.42, 12.0, 5.4, bullets, 17, MGRAY, b_font, 14, 8, "\u25B8")
                _page_num(s, sn, total, DGRAY)

            else:
                hdr_light(sec_cat)
                _rect(s, 0, 0.06, 0.06, 7.44, GOLD if idx % 2 == 0 else SEC)
                _bullet_block(s, 0.7, 1.42, 12.0, 5.4, bullets, 17, TBODY, b_font, 14, 8, "\u25B8")
                _page_num(s, sn, total, MGRAY)

    # ── CLOSING SLIDE ────────────────────────────────────────────
    if add_closing_slide:
        sn += 1
        s = prs.slides.add_slide(blank)
        set_bg(s, DARK)
        _rect(s, 0, 0, 13.333, 0.06, SEC)
        _rect(s, 0, 7.0, 13.333, 0.5, DARK2)
        # Decorative circle
        circ = s.shapes.add_shape(9, Inches(8.5), Inches(0.5), Inches(5.0), Inches(5.0))
        circ.fill.solid()
        circ.fill.fore_color.rgb = DARK2
        circ.line.fill.background()
        _rect(s, 5.4, 3.5, 2.5, 0.06, GOLD)
        _textbox(s, 0.8, 1.8, 7.5, 1.2, "Thank You", 48, WHITE, True, PP_ALIGN.LEFT, h_font)
        _textbox(s, 0.8, 3.8, 7.5, 0.55,
                 f"\u00A9{dt.datetime.now().year} {brand_label} | All rights reserved",
                 14, MGRAY, False, PP_ALIGN.LEFT, b_font)
        if brand_domain:
            _textbox(s, 0.8, 4.4, 7.5, 0.42, brand_domain, 13, DGRAY, False, font=b_font)
        _page_num(s, sn, total, DGRAY)

    name = f"{'report' if document_type in ('report', 'pitch_deck') else 'presentation'}_{uuid.uuid4().hex[:8]}.pptx"
    path = DATA_DIR / name
    prs.save(str(path))
    logger.info("Generated PPTX: %s (%d slides)", path, len(prs.slides))
    return path


def _parse_markdown_sections(content: str) -> list[dict[str, Any]]:
    return _parse_sections(content)
