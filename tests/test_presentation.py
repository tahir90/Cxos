"""Tests for the unified PPT generator (presentation.generate_pptx)."""

from __future__ import annotations

import shutil
from pathlib import Path
from unittest.mock import patch

import pytest
from pptx import Presentation

from agentic_cxo.tools.presentation import (
    _parse_markdown_sections,
    generate_pptx,
)

SAMPLE_REPORT = """\
# Executive Briefing: Cash-Flow Guardian Scenario

## Key Findings

1. **Top 3 Non-Payroll Expense Spikes:**
   - Vendor A: Increase of $50,000
   - Vendor B: Increase of $35,000
   - Vendor C: Increase of $25,000

2. **Duplicate Billing Audit:**
   - Vendor B: Identified $5,000 in duplicate charges.

## Risk Assessment

- **Documentation Risk:** Moderate
- **Cash reserves:** $8.2M with 5.9 months runway

## Recommendations

1. **Immediate:** Audit contractor contracts — potential $20-40k savings
2. **This week:** Consolidate SaaS tools
3. **This month:** Renegotiate top vendor contract

## Next Steps

- Schedule vendor review meeting
- Implement spend monitoring dashboard
"""


@pytest.fixture(autouse=True)
def clean_data():
    yield
    data_dir = Path(".cxo_data")
    if data_dir.exists():
        shutil.rmtree(data_dir)


class TestMarkdownParsing:
    def test_parse_sections(self):
        sections = _parse_markdown_sections(SAMPLE_REPORT)
        assert len(sections) >= 4
        titles = [s["title"] for s in sections]
        assert "Key Findings" in titles
        assert "Risk Assessment" in titles
        assert "Recommendations" in titles

    def test_parse_empty(self):
        sections = _parse_markdown_sections("")
        assert len(sections) == 1
        assert sections[0]["title"] == ""


class TestPptxMarkdownReport:
    def test_generates_pptx_file(self, tmp_path):
        with patch("agentic_cxo.tools.presentation.DATA_DIR", tmp_path):
            path = generate_pptx(
                SAMPLE_REPORT,
                title="The Cash-Flow Guardian",
                theme="dark",
                agent_role="CFO",
                sources=["quarterly_report.pdf", "vendor_contracts.pdf"],
                add_title_slide=True,
                add_closing_slide=True,
            )
        assert path.exists()
        assert path.suffix == ".pptx"

    def test_pptx_has_multiple_slides(self, tmp_path):
        with patch("agentic_cxo.tools.presentation.DATA_DIR", tmp_path):
            path = generate_pptx(
                SAMPLE_REPORT,
                title="Test Scenario",
                theme="dark",
                agent_role="CFO",
                add_title_slide=True,
                add_closing_slide=True,
            )
        prs = Presentation(str(path))
        assert len(prs.slides) >= 2

    def test_pptx_title_slide_content(self, tmp_path):
        with patch("agentic_cxo.tools.presentation.DATA_DIR", tmp_path):
            path = generate_pptx(
                SAMPLE_REPORT,
                title="The Tax Strategist",
                theme="dark",
                agent_role="CFO",
                add_title_slide=True,
                add_closing_slide=True,
            )
        prs = Presentation(str(path))
        first_slide = prs.slides[0]
        texts = [shape.text for shape in first_slide.shapes if shape.has_text_frame]
        combined = " ".join(texts)
        assert "Tax Strategist" in combined
        assert "CFO" in combined

    def test_pptx_no_sources(self, tmp_path):
        with patch("agentic_cxo.tools.presentation.DATA_DIR", tmp_path):
            path = generate_pptx(
                "## Finding\n\nSome content here.",
                title="Minimal",
                theme="dark",
                agent_role="COO",
                add_title_slide=True,
                add_closing_slide=True,
            )
        assert path.exists()
        prs = Presentation(str(path))
        assert len(prs.slides) >= 2

    def test_pptx_file_created_and_valid(self, tmp_path):
        with patch("agentic_cxo.tools.presentation.DATA_DIR", tmp_path):
            path = generate_pptx(
                SAMPLE_REPORT,
                title="The Cash-Flow Guardian",
                theme="dark",
                agent_role="CFO",
                sources=["quarterly_report.pdf"],
                add_title_slide=True,
                add_closing_slide=True,
            )
        assert path is not None
        assert path.exists()
        assert path.stat().st_size > 0
        prs = Presentation(str(path))
        assert len(prs.slides) >= 2
        first_slide = prs.slides[0]
        texts = [s.text for s in first_slide.shapes if s.has_text_frame]
        assert any("Cash-Flow Guardian" in t for t in texts)


class TestPptxOutline:
    def test_outline_light_theme(self):
        outline = "## Overview\n- Point A\n- Point B\n\n## Next Steps\n- Action 1"
        path = generate_pptx(outline, title="Test Deck", theme="light")
        assert path.exists()
        assert path.suffix == ".pptx"
        prs = Presentation(str(path))
        assert len(prs.slides) >= 2
