#!/usr/bin/env python3
"""Inspect PPTX structure and content for comparison."""

import sys
from pathlib import Path

def main():
    ppt_path = sys.argv[1] if len(sys.argv) > 1 else ".cxo_data/presentations/presentation_8e034d0a.pptx"
    p = Path(ppt_path)
    if not p.exists():
        p = Path("/workspace") / ppt_path
    if not p.exists():
        print("File not found:", ppt_path)
        return 1

    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(p))
    print(f"=== PPT INSPECTION: {p.name} ({p.stat().st_size} bytes) ===\n")
    print(f"Total slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        print(f"--- SLIDE {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        font_size = para.font.size.pt if para.font.size else 0
                        print(f"  [{font_size:.0f}pt] {text[:120]}")
            if shape.shape_type == 14:  # placeholder
                pass
        print()
    return 0

if __name__ == "__main__":
    sys.exit(main())
